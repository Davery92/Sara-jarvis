import os
import mimetypes
import logging
from typing import List, Dict, Any, Optional, Tuple
from io import BytesIO
import tempfile

# Text extraction libraries
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from readability import Document as ReadabilityDocument
except ImportError:
    ReadabilityDocument = None

from markdownify import markdownify
from minio import Minio
from minio.error import S3Error

from app.core.config import settings
from app.services.embeddings import chunk_text, get_embeddings_batch

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing and extracting text from documents"""
    
    def __init__(self):
        # Initialize MinIO client
        try:
            self.minio_client = Minio(
                settings.minio_url.replace("http://", ""),
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=False
            )
            
            # Ensure bucket exists
            if not self.minio_client.bucket_exists(settings.minio_bucket):
                self.minio_client.make_bucket(settings.minio_bucket)
                logger.info(f"Created MinIO bucket: {settings.minio_bucket}")
                
        except Exception as e:
            logger.error(f"Failed to initialize MinIO client: {e}")
            self.minio_client = None
    
    def extract_text(self, file_content: bytes, mime_type: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from various file formats"""
        
        text = ""
        metadata = {"filename": filename, "mime_type": mime_type}
        
        try:
            if mime_type == "application/pdf":
                text, pdf_meta = self._extract_pdf(file_content)
                metadata.update(pdf_meta)
                
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text, docx_meta = self._extract_docx(file_content)
                metadata.update(docx_meta)
                
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text, pptx_meta = self._extract_pptx(file_content)
                metadata.update(pptx_meta)
                
            elif mime_type.startswith("text/"):
                text = file_content.decode("utf-8", errors="ignore")
                metadata["encoding"] = "utf-8"
                
            elif mime_type == "text/html":
                text = self._extract_html(file_content)
                metadata["extraction_method"] = "readability"
                
            elif mime_type == "text/markdown":
                text = file_content.decode("utf-8", errors="ignore")
                metadata["format"] = "markdown"
                
            else:
                # Try to decode as text
                try:
                    text = file_content.decode("utf-8", errors="ignore")
                    metadata["extraction_method"] = "utf-8_fallback"
                except:
                    raise ValueError(f"Unsupported file type: {mime_type}")
            
            metadata["text_length"] = len(text)
            metadata["extraction_success"] = True
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Text extraction failed for {filename}: {e}")
            metadata["extraction_error"] = str(e)
            metadata["extraction_success"] = False
            return "", metadata
    
    def _extract_pdf(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF"""
        if not pypdf:
            raise ImportError("pypdf not available for PDF extraction")
        
        text_parts = []
        metadata = {"pages": 0}
        
        try:
            with BytesIO(content) as pdf_file:
                reader = pypdf.PdfReader(pdf_file)
                metadata["pages"] = len(reader.pages)
                
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                
                # Extract metadata if available
                if reader.metadata:
                    if reader.metadata.title:
                        metadata["title"] = reader.metadata.title
                    if reader.metadata.author:
                        metadata["author"] = reader.metadata.author
                    if reader.metadata.subject:
                        metadata["subject"] = reader.metadata.subject
                        
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise
        
        return "\n\n".join(text_parts), metadata
    
    def _extract_docx(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX"""
        if not DocxDocument:
            raise ImportError("python-docx not available for DOCX extraction")
        
        text_parts = []
        metadata = {"paragraphs": 0, "tables": 0}
        
        try:
            with BytesIO(content) as docx_file:
                doc = DocxDocument(docx_file)
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                        metadata["paragraphs"] += 1
                
                # Extract table content
                for table in doc.tables:
                    metadata["tables"] += 1
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_parts.append(" | ".join(row_text))
                            
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise
        
        return "\n\n".join(text_parts), metadata
    
    def _extract_pptx(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX"""
        if not Presentation:
            raise ImportError("python-pptx not available for PPTX extraction")
        
        text_parts = []
        metadata = {"slides": 0}
        
        try:
            with BytesIO(content) as pptx_file:
                prs = Presentation(pptx_file)
                metadata["slides"] = len(prs.slides)
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    slide_text.append(f"[Slide {slide_num + 1}]")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    
                    if len(slide_text) > 1:  # More than just the slide number
                        text_parts.append("\n".join(slide_text))
                        
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise
        
        return "\n\n".join(text_parts), metadata
    
    def _extract_html(self, content: bytes) -> str:
        """Extract text from HTML using readability"""
        if not ReadabilityDocument:
            # Fallback to basic HTML stripping
            html = content.decode("utf-8", errors="ignore")
            # Basic HTML tag removal (very simple)
            import re
            return re.sub(r'<[^>]+>', '', html)
        
        try:
            html = content.decode("utf-8", errors="ignore")
            doc = ReadabilityDocument(html)
            
            # Get the main content and convert to markdown
            content_html = doc.summary()
            text = markdownify(content_html, heading_style="ATX")
            
            return text
            
        except Exception as e:
            logger.error(f"HTML extraction error: {e}")
            # Fallback to basic stripping
            html = content.decode("utf-8", errors="ignore")
            import re
            return re.sub(r'<[^>]+>', '', html)
    
    async def store_file(self, file_content: bytes, filename: str, mime_type: str) -> str:
        """Store file in MinIO and return storage key"""
        
        if not self.minio_client:
            raise RuntimeError("MinIO client not available")
        
        # Generate unique storage key
        import uuid
        file_ext = os.path.splitext(filename)[1]
        storage_key = f"{uuid.uuid4()}{file_ext}"
        
        try:
            # Upload to MinIO
            self.minio_client.put_object(
                bucket_name=settings.minio_bucket,
                object_name=storage_key,
                data=BytesIO(file_content),
                length=len(file_content),
                content_type=mime_type
            )
            
            logger.info(f"Stored file {filename} as {storage_key}")
            return storage_key
            
        except S3Error as e:
            logger.error(f"MinIO storage error: {e}")
            raise RuntimeError(f"Failed to store file: {e}")
    
    def get_file(self, storage_key: str) -> bytes:
        """Retrieve file from MinIO"""
        
        if not self.minio_client:
            raise RuntimeError("MinIO client not available")
        
        try:
            response = self.minio_client.get_object(
                bucket_name=settings.minio_bucket,
                object_name=storage_key
            )
            
            content = response.read()
            response.close()
            response.release_conn()
            
            return content
            
        except S3Error as e:
            logger.error(f"MinIO retrieval error: {e}")
            raise RuntimeError(f"Failed to retrieve file: {e}")
    
    def delete_file(self, storage_key: str) -> bool:
        """Delete file from MinIO"""
        
        if not self.minio_client:
            return False
        
        try:
            self.minio_client.remove_object(
                bucket_name=settings.minio_bucket,
                object_name=storage_key
            )
            
            logger.info(f"Deleted file {storage_key}")
            return True
            
        except S3Error as e:
            logger.error(f"MinIO deletion error: {e}")
            return False
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str, 
        mime_type: str = None
    ) -> Tuple[str, List[str], Dict[str, Any], str]:
        """
        Complete document processing pipeline
        Returns: (extracted_text, chunks, metadata, storage_key)
        """
        
        # Detect MIME type if not provided
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)
            mime_type = mime_type or "application/octet-stream"
        
        # Store file
        storage_key = await self.store_file(file_content, filename, mime_type)
        
        # Extract text
        text, metadata = self.extract_text(file_content, mime_type, filename)
        
        if not text.strip():
            raise ValueError("No text could be extracted from the document")
        
        # Create chunks
        chunks = chunk_text(
            text,
            chunk_size=settings.memory_chunk_size,
            overlap=settings.memory_chunk_overlap
        )
        
        metadata.update({
            "storage_key": storage_key,
            "chunk_count": len(chunks),
            "processing_success": True
        })
        
        return text, chunks, metadata, storage_key


# Global document processor instance
doc_processor = DocumentProcessor()